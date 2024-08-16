from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title_1 = slide_1.shapes.title
subtitle_1 = slide_1.placeholders[1]

title_1.text = "Elevator Pitch Script for Cloud DevOps Platform"
subtitle_1.text = "AccelSDLC by Saravanan Gnanaguru"

# Slide 2: Opening Hook
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title_2 = slide_2.shapes.title
content_2 = slide_2.shapes.placeholders[1]

title_2.text = "Opening Hook"
content_2.text = (
    "Imagine you're the CTO of a product development company, working with a team of developers who find themselves stuck when "
    "application issues arise after deploying to production. The stress of troubleshooting in real-time is not only frustrating but "
    "also disrupts your team's workflow.\n\n"
    "So, what if there was an Automation Platform that your developers could use to seamlessly onboard applications and new features, "
    "ensuring everything works perfectly from the start? That’s exactly what I’m ideating right now."
)

# Slide 3: Problem Statement
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title_3 = slide_3.shapes.title
content_3 = slide_3.shapes.placeholders[1]

title_3.text = "Problem Statement"
content_3.text = (
    "In today’s competitive product development landscape, even the most skilled developers can struggle with the complexities of deploying and "
    "maintaining applications in production. Manual processes, unexpected errors, and the lack of a unified system for monitoring and managing both "
    "applications and infrastructure lead to delays and inefficiencies, ultimately impacting your bottom line."
)

# Slide 4: Solution Introduction
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title_4 = slide_4.shapes.title
content_4 = slide_4.shapes.placeholders[1]

title_4.text = "Solution Introduction"
content_4.text = (
    "Introducing our PaaS product offering AccelSDLC — a Cloud DevOps platform designed to eliminate these challenges by automating the entire software "
    "development lifecycle. From onboarding new applications to deploying features and monitoring performance, AccelSDLC ensures that your developers can "
    "focus on writing code, while our platform takes care of the rest.\n\n"
    "We're also planning to integrate AccelSDLC with major cloud providers and popular automation tools, allowing teams to choose customized solutions that "
    "best fit their needs. Additionally, we have plans to create our own cloud offering that will integrate seamlessly with our platform, providing an "
    "all-in-one solution for development and deployment.\n\n"
    "This product has potential for micro SaaS offerings of sub features, like containerisation of apps using template, CICD automation for specific applications."
)

# Slide 5: Value Proposition
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title_5 = slide_5.shapes.title
content_5 = slide_5.shapes.placeholders[1]

title_5.text = "Value Proposition"
content_5.text = (
    "With AccelSDLC, your team gains access to an integrated platform that not only automates deployment but also provides comprehensive monitoring for both "
    "applications and infrastructure. Real-time dashboards and detailed reporting give you full visibility into every aspect of your development and operations, "
    "helping you identify and resolve issues before they affect your users. Our platform engineering features make it easy for developers to onboard new applications "
    "and features with minimal effort, while our future cloud integration will offer even more flexibility and control."
)

# Slide 6: Market Validation
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title_6 = slide_6.shapes.title
content_6 = slide_6.shapes.placeholders[1]

title_6.text = "Market Validation"
content_6.text = (
    "The demand for DevOps and cloud platforms is surging, with the market expected to grow from $7 billion in 2020 to over $17 billion by 2026. "
    "Companies adopting DevOps practices see a 22% increase in software delivery performance and a 50% reduction in time-to-market. "
    "By leveraging AccelSDLC, you can stay ahead of this trend, ensuring your development processes are not just up-to-date but leading the industry."
)

# Slide 7: Call to Action
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
title_7 = slide_7.shapes.title
content_7 = slide_7.shapes.placeholders[1]

title_7.text = "Call to Action"
content_7.text = (
    "I’d love to connect with you for a detailed discussion on how AccelSDLC can transform your development process. Let's schedule a call to dive into "
    "the unique selling propositions and how our platform can be tailored to meet your specific needs. Please reach out to me directly, or visit "
    "www.accelsdlc.com to set up a time that works best for you."
)

# Slide 8: Closing Statement
slide_8 = prs.slides.add_slide(prs.slide_layouts[1])
title_8 = slide_8.shapes.title
content_8 = slide_8.shapes.placeholders[1]

title_8.text = "Closing Statement"
content_8.text = (
    "Thank you for your time. I’m Saravanan Gnanaguru, and I’m excited to explore how AccelSDLC can elevate your product development to new heights."
)

# Save the presentation
pptx_path = "/Users/gsaravanan/Documents/AccelSDLC_Pitch_Deck.pptx"
prs.save(pptx_path)

pptx_path
